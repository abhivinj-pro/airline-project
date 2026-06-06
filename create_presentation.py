"""
Generate AirWave AI Booking Recovery MVP presentation for CTO audience.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x00, 0xD4, 0xFF)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
ACCENT3 = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MID_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tx_box


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT, spacing=Pt(8)):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"
        p.space_after = spacing
    return tx_box


def add_icon_card(slide, left, top, width, height, icon_text, title, description,
                  icon_color=ACCENT, bg_color=BG_CARD):
    card = add_shape_bg(slide, left, top, width, height, bg_color, corner_radius=0.05)
    circle_size = Inches(0.7)
    cx = left + Inches(0.3)
    cy = top + Inches(0.3)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, circle_size, circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_color
    circle.line.fill.background()

    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, left + Inches(0.3), top + Inches(1.2), width - Inches(0.6), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.7), width - Inches(0.6), height - Inches(2.0),
                 description, font_size=13, color=LIGHT_GRAY)
    return card


def add_metric_box(slide, left, top, width, height, value, label, color=ACCENT):
    card = add_shape_bg(slide, left, top, width, height, BG_CARD, corner_radius=0.05)
    add_text_box(slide, left, top + Inches(0.25), width, Inches(0.6),
                 value, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.9), width, Inches(0.5),
                 label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return card


def add_divider_line(slide, left, top, width, color=ACCENT):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# Slide 1: Title
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "AirWave AI", font_size=60, color=ACCENT, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.8),
             "Booking Recovery System", font_size=44, color=WHITE, bold=True)
add_divider_line(slide, Inches(1.5), Inches(3.9), Inches(3), ACCENT)
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(9), Inches(0.7),
             "AI-powered abandonment detection and real-time recovery for airline direct channels",
             font_size=20, color=LIGHT_GRAY)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(5), Inches(0.5),
             "MVP Technical Overview  |  April 2026", font_size=16, color=LIGHT_GRAY)


# Slide 2: Problem
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "THE PROBLEM", font_size=14, color=RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Airlines Are Bleeding Revenue at Checkout", font_size=36, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.7), Inches(2.5), RED)

add_metric_box(slide, Inches(0.8), Inches(2.2), Inches(3.5), Inches(1.5),
               "50%", "Booking sessions abandoned before payment", RED)
add_metric_box(slide, Inches(4.8), Inches(2.2), Inches(3.5), Inches(1.5),
               "50,000", "Abandoned sessions / month (per 100K starts)", ORANGE)
add_metric_box(slide, Inches(8.8), Inches(2.2), Inches(3.5), Inches(1.5),
               "INR 60Cr", "Annual gross booking value at risk", RED)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
             "What happens when users abandon:", font_size=20, color=WHITE, bold=True)
problem_items = [
    "Lost direct revenue — users book through OTAs or competitors instead",
    "High commission costs — offline agents and aggregators capture the booking",
    "Missed ancillary revenue — no cross-sell of seats, bags, meals, or insurance",
    "No visibility into why users leave — no data to improve the funnel",
]
add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(2.5),
                problem_items, font_size=17, color=MID_GRAY, bullet_color=RED)


# Slide 3: Solution overview
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "THE SOLUTION", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "3-Pillar AI Recovery System", font_size=36, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.7), Inches(2.5), ACCENT)

card_w = Inches(3.6)
card_h = Inches(4.5)
gap = Inches(0.45)
start_x = Inches(0.8)
card_y = Inches(2.2)

add_icon_card(slide, start_x, card_y, card_w, card_h,
              "1", "Behavioral Abandonment Detection",
              "Client-side tracker monitors 10+ real-time signals — idle time, back button, tab switches, rapid clicks, price hovering, form fatigue — weighted by booking step. Composite risk score triggers tiered interventions.",
              icon_color=ACCENT)

add_icon_card(slide, start_x + card_w + gap, card_y, card_w, card_h,
              "2", "Proactive AI Chat Concierge",
              "Azure OpenAI-powered assistant engages users when risk score crosses thresholds. Context-aware: knows the flight, price, funnel step, and specific trigger signals. Adapts tone from friendly nudge (0.45+) to urgent save (0.80+).",
              icon_color=ACCENT2)

add_icon_card(slide, start_x + 2 * (card_w + gap), card_y, card_w, card_h,
              "3", "Smart Recovery Campaigns",
              "For users who already left: AI generates personalized email/SMS with flight details, seat scarcity, and deep links back to their saved session. Users resume exactly where they stopped — all details pre-filled.",
              icon_color=ACCENT3)


# Slide 4: Detection engine detail
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "PILLAR 1", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Real-Time Behavioral Abandonment Detection", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ACCENT)

signals = [
    ("idle_60s", "0.35", "User distracted / indecisive"),
    ("back_button", "0.25", "Reconsidering their choice"),
    ("idle_30s", "0.20", "Hesitation signal"),
    ("mouse_leave", "0.15", "About to leave the site"),
    ("tab_switch", "0.12", "Comparison shopping"),
    ("form_delete", "0.10", "Form fatigue"),
    ("scroll_top", "0.08", "Reconsidering journey"),
    ("rapid_clicks", "0.06", "UI frustration"),
    ("price_hover", "0.05", "Price sensitivity"),
]

bar_start_x = Inches(0.8)
bar_y = Inches(2.1)
max_bar_w = Inches(4.0)

for i, (signal, weight, _desc) in enumerate(signals):
    y = bar_y + Inches(i * 0.52)
    w_frac = float(weight) / 0.35
    bar_w = max(Inches(0.3), int(max_bar_w * w_frac))

    bar_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        bar_start_x + Inches(1.8),
        y,
        bar_w,
        Inches(0.35),
    )
    bar_shape.fill.solid()
    r = int(0x00 + (0x7C - 0x00) * i / 8)
    g = int(0xD4 + (0x3A - 0xD4) * i / 8)
    b = int(0xFF + (0xED - 0xFF) * i / 8)
    bar_shape.fill.fore_color.rgb = RGBColor(r, g, b)
    bar_shape.line.fill.background()
    bar_shape.adjustments[0] = 0.15

    add_text_box(slide, bar_start_x, y, Inches(1.7), Inches(0.35),
                 signal, font_size=12, color=ACCENT, bold=True)
    tf = bar_shape.text_frame
    p = tf.paragraphs[0]
    p.text = weight
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"

right_x = Inches(7.2)
add_text_box(slide, right_x, Inches(2.1), Inches(5.5), Inches(0.6),
             "How It Works", font_size=22, color=WHITE, bold=True)
explanations = [
    "Signals are weighted and combined into a composite risk score",
    "Booking step acts as a multiplier — later steps mean higher intent and higher recovery priority",
    "Three intervention thresholds: 0.45 (nudge), 0.65 (assist), 0.80 (urgent save)",
    "MVP uses heuristic weights; production would use a supervised ML model trained on historical funnel data",
]
add_bullet_list(slide, right_x, Inches(2.7), Inches(5.5), Inches(3.5),
                explanations, font_size=15, color=MID_GRAY, bullet_color=ACCENT,
                spacing=Pt(12))


# Slide 5: AI concierge detail
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "PILLAR 2", font_size=14, color=ACCENT2, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Proactive AI Chat Concierge", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ACCENT2)

tiers = [
    ("Score 0.45+", "Friendly Nudge", '"Need help finding the right flight?"',
     "Low-pressure engagement to offer assistance", ACCENT3),
    ("Score 0.65+", "Active Assist", '"I noticed you\'re on passenger details — can I help?"',
     "Context-specific help based on funnel step", ORANGE),
    ("Score 0.80+", "Urgent Save", '"Your booking is saved. Want me to help complete it?"',
     "Maximum effort to retain high-intent user", RED),
]

for i, (score, title, quote, desc, color) in enumerate(tiers):
    y = Inches(2.2) + Inches(i * 1.55)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(1.8), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.adjustments[0] = 0.3
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = score
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(3.0), y - Inches(0.05), Inches(4), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.0), y + Inches(0.35), Inches(5), Inches(0.4),
                 quote, font_size=15, color=ACCENT)
    add_text_box(slide, Inches(3.0), y + Inches(0.75), Inches(5), Inches(0.4),
                 desc, font_size=13, color=LIGHT_GRAY)

right_x = Inches(8.5)
add_text_box(slide, right_x, Inches(2.2), Inches(4), Inches(0.5),
             "AI Context Awareness", font_size=20, color=WHITE, bold=True)
ctx_items = [
    "Flight route, airline, and price selected",
    "Current booking funnel step",
    "Which behavioral signals triggered",
    "Price sensitivity vs. form fatigue vs. comparison shopping",
    "Adapts tone and strategy per user context",
]
add_bullet_list(slide, right_x, Inches(2.8), Inches(4.2), Inches(3.5),
                ctx_items, font_size=14, color=MID_GRAY, bullet_color=ACCENT2,
                spacing=Pt(10))


# Slide 6: Recovery campaigns
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "PILLAR 3", font_size=14, color=ACCENT3, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Smart Recovery Campaigns", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ACCENT3)

add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), BG_CARD, 0.03)
add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.4),
             "EMAIL RECOVERY", font_size=12, color=ACCENT3, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.9), Inches(5), Inches(0.5),
             "Sample Personalized Recovery Email", font_size=18, color=WHITE, bold=True)

email_lines = [
    '"Your Dubai flight at INR 32,000 is still available"',
    "",
    ">  DEL → DXB  |  AI-205  |  Mar 15, 2026",
    "$  INR 32,000  |  Only 3 seats left",
    "",
    "Deep link back to saved session",
    "All passenger details pre-filled",
    "One-click resume to payment step",
]
email_y = Inches(3.5)
for line in email_lines:
    if not line:
        email_y += Inches(0.2)
        continue
    color = ACCENT if line.startswith('"') else MID_GRAY
    size = 15 if line.startswith('"') else 13
    add_text_box(slide, Inches(1.3), email_y, Inches(4.8), Inches(0.35),
                 line, font_size=size, color=color)
    email_y += Inches(0.35)

right_x = Inches(7.0)
add_text_box(slide, right_x, Inches(2.2), Inches(5.5), Inches(0.5),
             "Campaign Intelligence", font_size=22, color=WHITE, bold=True)
camp_items = [
    "AI generates each message from the specific customer journey",
    "Personalized with route, price, scarcity, and timing signals",
    "Multi-channel: email for detail-rich, SMS for urgency",
    "Deep links restore full session state — zero re-entry friction",
    "Targets only high-intent, contactable abandoners",
]
add_bullet_list(slide, right_x, Inches(2.8), Inches(5.5), Inches(3.0),
                camp_items, font_size=15, color=MID_GRAY, bullet_color=ACCENT3,
                spacing=Pt(10))


# Slide 7: Future capabilities
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "ROADMAP", font_size=14, color=ACCENT2, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "AI-Powered Retention: Beyond the MVP", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ACCENT2)

future = [
    ("Abandonment Reason\nClassification",
     "AI classifies why users leave — price, schedule, policy, trust, payment, or delay — so interventions match the real objection, not just the risk score.",
     ACCENT, "Highest Impact"),
    ("AI Timing & Channel\nOptimizer",
     "Predicts optimal send time, cadence, and channel (email vs. SMS vs. chat) per abandoner. Fewer, better-timed interventions outperform blanket reminders.",
     ACCENT2, "High Impact"),
    ("AI-Assisted Human\nSave Desk",
     "For high-value bookings: AI summarizes session, predicts intent, surfaces objections, and recommends talking points — so human agents recover faster.",
     ACCENT3, "Strategic"),
]

for i, (title, desc, color, badge_text) in enumerate(future):
    x = Inches(0.8) + Inches(i * 4.05)
    add_shape_bg(slide, x, Inches(2.3), Inches(3.7), Inches(4.5), BG_CARD, 0.04)

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(0.25), Inches(2.5), Inches(1.6), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.adjustments[0] = 0.3
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = badge_text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.25), Inches(3.1), Inches(3.2), Inches(0.9),
                 title, font_size=19, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.25), Inches(3.95), Inches(3.2), Inches(2.5),
                 desc, font_size=13, color=LIGHT_GRAY)


# Slide 8: Price lock impact
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "BUSINESS CASE", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Estimated Impact of a 4-Hour Price Lock", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ACCENT)

price_lock_metrics = [
    ("10,000", "Strong-candidate\nabandoners / month", ACCENT),
    ("800", "Users accepting or\nresponding to offer", ORANGE),
    ("96", "Incremental recovered\nbookings", ACCENT3),
    ("INR 1.13M", "Net expected monthly\nrevenue change", WHITE),
]
for i, (value, label, color) in enumerate(price_lock_metrics):
    x = Inches(0.8) + Inches(i * 3.05)
    add_metric_box(slide, x, Inches(2.0), Inches(2.7), Inches(1.3), value, label, color)

add_shape_bg(slide, Inches(0.8), Inches(3.65), Inches(5.1), Inches(2.3), BG_CARD, 0.04)
add_text_box(slide, Inches(1.1), Inches(3.85), Inches(4.7), Inches(0.4),
             "Assumptions", font_size=20, color=WHITE, bold=True)
price_lock_items = [
    "20% of abandoners qualify for a 4-hour fare hold → 10,000 users/month",
    "8% accept or meaningfully respond → 800 treated users",
    "Baseline conversion for this subgroup = 20% → 160 bookings",
    "Price lock adds 12 percentage points → 32% treated conversion",
    "Average booking value = INR 12,000",
]
add_bullet_list(slide, Inches(1.1), Inches(4.25), Inches(4.6), Inches(1.45),
                price_lock_items, font_size=13, color=MID_GRAY, bullet_color=ACCENT,
                spacing=Pt(5))

add_shape_bg(slide, Inches(6.2), Inches(3.65), Inches(6.1), Inches(2.3), BG_CARD, 0.04)
add_text_box(slide, Inches(6.5), Inches(3.85), Inches(5.5), Inches(0.4),
             "Expected-Value Calculation", font_size=20, color=WHITE, bold=True)
calc_lines = [
    "E(RevenueChange) = Revenue increase from recovered bookings - Locked-fare downside",
    "Recovered revenue: 96 x INR 12,000 = INR 1,152,000",
    "Fare-lock downside: 160 x (15% x INR 1,000) = INR 24,000",
    "Net monthly expected value: INR 1,152,000 - INR 24,000 = INR 1,128,000",
]
calc_colors = [WHITE, MID_GRAY, MID_GRAY, ACCENT3]
for i, line in enumerate(calc_lines):
    add_text_box(slide, Inches(6.5), Inches(4.28) + Inches(i * 0.38), Inches(5.4), Inches(0.3),
                 line, font_size=12 if i == 0 else 13, color=calc_colors[i], bold=(i == 3))

add_shape_bg(slide, Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.75), BG_CARD, 0.04)
add_text_box(slide, Inches(1.0), Inches(6.34), Inches(11.1), Inches(0.35),
             "Interpretation: under these assumptions, the 4-hour price lock has a strongly positive expected value. The upside from 96 incremental recovered bookings materially outweighs the expected fare movement given up on bookings that would likely have happened anyway.",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)


# Slide 9: Campaign impact
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "BUSINESS CASE", font_size=14, color=ACCENT3, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Estimated Impact of Booking-Recovery Campaigns", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ACCENT3)

campaign_metrics = [
    ("15,000", "High-intent and\ncontactable users / month", ACCENT3),
    ("3%", "Recovery rate across\nemail, SMS, and chat", ORANGE),
    ("450", "Incremental recovered\nbookings", ACCENT),
    ("INR 5.40M", "Estimated gross booking\nvalue recovered / month", WHITE),
]
for i, (value, label, color) in enumerate(campaign_metrics):
    x = Inches(0.8) + Inches(i * 3.05)
    add_metric_box(slide, x, Inches(2.0), Inches(2.7), Inches(1.3), value, label, color)

add_shape_bg(slide, Inches(0.8), Inches(3.65), Inches(5.1), Inches(2.3), BG_CARD, 0.04)
add_text_box(slide, Inches(1.1), Inches(3.85), Inches(4.6), Inches(0.4),
             "Assumptions", font_size=20, color=WHITE, bold=True)
campaign_items = [
    "30% of abandoned sessions are high-intent and contactable → 15,000 users/month",
    "Email, SMS, and chat journeys recover 3% of that audience",
    "Recovered bookings: 15,000 x 3% = 450",
    "Average booking value = INR 12,000",
    "Recovered gross booking value: 450 x INR 12,000 = INR 5,400,000",
]
add_bullet_list(slide, Inches(1.1), Inches(4.25), Inches(4.6), Inches(1.45),
                campaign_items, font_size=13, color=MID_GRAY, bullet_color=ACCENT3,
                spacing=Pt(5))

add_shape_bg(slide, Inches(6.2), Inches(3.65), Inches(6.1), Inches(2.3), BG_CARD, 0.04)
add_text_box(slide, Inches(6.5), Inches(3.85), Inches(5.5), Inches(0.4),
             "Why Campaigns Scale", font_size=20, color=WHITE, bold=True)
campaign_scale_items = [
    "Campaigns reach a larger portion of abandoners than fare-lock offers",
    "They address distraction, form fatigue, comparison shopping, payment hesitation, and delayed decisions",
    "They compound with saved sessions, proactive chat, and deep-link resume flows",
]
add_bullet_list(slide, Inches(6.5), Inches(4.25), Inches(5.3), Inches(1.45),
                campaign_scale_items, font_size=13, color=MID_GRAY, bullet_color=ACCENT,
                spacing=Pt(6))

add_shape_bg(slide, Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.75), BG_CARD, 0.04)
add_text_box(slide, Inches(1.0), Inches(6.34), Inches(11.1), Inches(0.35),
             "Interpretation: campaigns produce materially larger revenue impact because they address a broader set of abandonment causes across a much larger reachable audience.",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)


# Slide 10: Comparison
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "BUSINESS IMPACT", font_size=14, color=ORANGE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Comparison: Price Lock vs. Recovery Campaigns", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ORANGE)

comparison_metrics = [
    ("INR 1.13M", "4-hour price lock\nnet expected value / month", ACCENT),
    ("INR 5.40M", "Recovery campaigns\nvalue recovered / month", ACCENT3),
    ("4.8x", "Campaign impact vs.\nprice lock", ORANGE),
]
for i, (value, label, color) in enumerate(comparison_metrics):
    x = Inches(0.8) + Inches(i * 4.05)
    add_metric_box(slide, x, Inches(2.0), Inches(3.6), Inches(1.4), value, label, color)

add_shape_bg(slide, Inches(0.8), Inches(3.9), Inches(5.4), Inches(2.5), BG_CARD, 0.04)
add_text_box(slide, Inches(1.1), Inches(4.1), Inches(4.8), Inches(0.4),
             "Why the Gap Is Reasonable", font_size=20, color=WHITE, bold=True)
gap_items = [
    "Price lock is relevant only for users whose main barrier is fare volatility",
    "Campaigns address a much larger pool of abandonment causes across the funnel",
    "Price lock is best treated as one tactic inside the larger recovery system",
]
add_bullet_list(slide, Inches(1.1), Inches(4.55), Inches(4.8), Inches(1.4),
                gap_items, font_size=13, color=MID_GRAY, bullet_color=ORANGE,
                spacing=Pt(6))

add_shape_bg(slide, Inches(6.5), Inches(3.9), Inches(5.8), Inches(2.5), BG_CARD, 0.04)
add_text_box(slide, Inches(6.8), Inches(4.1), Inches(5.0), Inches(0.4),
             "Strategic Takeaway", font_size=20, color=WHITE, bold=True)
takeaway_lines = [
    "The larger revenue upside comes from recovering abandoned bookings at scale through proactive chat and post-abandonment outreach.",
    "Combined modeled opportunity: INR 6.53M/month, with campaigns contributing the majority of the upside.",
]
for i, line in enumerate(takeaway_lines):
    add_text_box(slide, Inches(6.8), Inches(4.65) + Inches(i * 0.55), Inches(5.0), Inches(0.45),
                 line, font_size=14, color=ACCENT3 if i == 0 else WHITE, bold=(i == 0))


# Slide 11: MVP vs production
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "TECHNICAL MATURITY", font_size=14, color=ORANGE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "MVP Assumptions vs. Production Path", font_size=32, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.65), Inches(2.5), ORANGE)

items_left = [
    ("Risk scoring", "Heuristic weights", "Supervised ML model on historical data"),
    ("Event signals", "10 browser events", "Rich event taxonomy + device/channel context"),
    ("Inventory", "Static demo data", "Live airline PSS/NDC API integration"),
    ("Identity", "Browser session only", "Server tokens + CRM/CDP identity stitching"),
    ("AI promises", "Mentions unverified actions", "Policy + tool-calling layer with verified APIs"),
]
items_right = [
    ("Abandonment def.", "Score threshold only", "Business rules: inactivity + step + payment + suppression"),
    ("Recovery delivery", "Simulated (SQLite log)", "ESP/SMS providers + delivery/open/click tracking"),
    ("Price lock", "Assumed 4hr hold", "Real fare-lock via pricing/inventory systems"),
    ("Checkout", "Simulated conversion", "Payment gateway + PNR creation + ticketing"),
    ("PII handling", "Local DB, no encryption", "Encrypted, tokenized, consent-aware, auditable"),
]

for col, items in enumerate([items_left, items_right]):
    x = Inches(0.8) if col == 0 else Inches(6.8)
    for i, (area, mvp, prod) in enumerate(items):
        y = Inches(2.1) + Inches(i * 1.02)
        add_text_box(slide, x, y, Inches(5.5), Inches(0.3),
                     area.upper(), font_size=11, color=ORANGE, bold=True)
        add_text_box(slide, x, y + Inches(0.28), Inches(2.6), Inches(0.3),
                     "MVP: " + mvp, font_size=12, color=LIGHT_GRAY)
        add_text_box(slide, x, y + Inches(0.55), Inches(5.5), Inches(0.35),
                     "→ Prod: " + prod, font_size=12, color=WHITE)


# Slide 12: Key takeaways
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "KEY TAKEAWAYS", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.8),
             "Why AirWave AI Matters", font_size=36, color=WHITE, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.7), Inches(2.5), ACCENT)

takeaways = [
    ("Real-Time Detection", "10+ behavioral signals scored in real-time to catch abandonment before it happens — not after.", ACCENT),
    ("AI-Powered, Context-Aware", "The concierge knows the flight, the price, the step, and the reason for hesitation. It adapts — not a scripted chatbot.", ACCENT2),
    ("Recover at Scale", "Personalized multi-channel campaigns with deep links that restore full session state. Zero-friction re-entry.", ACCENT3),
    ("Measurable Revenue Impact", "Estimated INR 6.5M+/month in recovered bookings. Recovery campaigns alone deliver about 4.8x more than price lock.", ORANGE),
]

for i, (title, desc, color) in enumerate(takeaways):
    y = Inches(2.2) + Inches(i * 1.2)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.7), y - Inches(0.05), Inches(10), Inches(0.4),
                 title, font_size=22, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.4), Inches(10), Inches(0.5),
                 desc, font_size=16, color=MID_GRAY)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             '"A 4-hour price lock is a useful conversion aid, but the larger revenue upside comes from recovering abandoned bookings at scale through proactive chat and post-abandonment outreach."',
             font_size=15, color=ACCENT, alignment=PP_ALIGN.CENTER)


output_path = Path(r"c:\repos\Random\airline-booking-recovery\AirWave_AI_MVP_Presentation.pptx")
try:
    prs.save(output_path)
    saved_path = output_path
except PermissionError:
    saved_path = output_path.with_name(output_path.stem + "_updated.pptx")
    prs.save(saved_path)

print(f"Presentation saved to: {saved_path}")
